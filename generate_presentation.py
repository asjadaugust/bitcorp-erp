from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    blank_slide_layout = prs.slide_layouts[6]
    
    # Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Bitcorp ERP - Modernización e Inteligencia Artificial"
    subtitle.text = "Casos de Uso: LLMs, Machine Learning, Visión Computacional \nGenerado con Nano Banana 2"
    
    images = [
        {
            "img": "/Users/klm95441/.gemini/antigravity/brain/c77360cc-e824-4c0d-b170-0c52ac0981c5/ai_erp_title_1772731397170.png",
            "title": "El Futuro del ERP para Gestión de Equipos",
            "desc": "Transformando el Procedimiento CORP-GEM-P-001 y CORP-GEM-P-002 con Inteligencia Artificial. Optimización de presupuestos, eficiencia operativa y prevención de riesgos."
        },
        {
            "img": "/Users/klm95441/.gemini/antigravity/brain/c77360cc-e824-4c0d-b170-0c52ac0981c5/chatbot_operator_1772731416441.png",
            "title": "Chatbot / LLM Asistente para Operadores",
            "desc": "Agilice el Parte Diario de Equipos y Vales de Combustible.\n- Los operadores reportan sus horas por WhatsApp.\n- Inteligencia Artificial extrae horómetros y galones de combustible.\n- Válido para CORP-GEM-F-005. Reducción en errores de tipeo y demoras en campo."
        },
        {
            "img": "/Users/klm95441/.gemini/antigravity/brain/c77360cc-e824-4c0d-b170-0c52ac0981c5/cv_excavator_1772731432647.png",
            "title": "Visión Computacional en Checklist",
            "desc": "Automatice la Inspección y Devolución de Maquinaria.\n- Análisis de fotografías para el formato CORP-SSOMA-F-073.\n- Detección de daños: abolladuras, fugas hidráulicas, estado de neumáticos.\n- Transparencia con el proveedor en 'Acta de Entrega' (CORP-GEM-F-010)."
        },
        {
            "img": "/Users/klm95441/.gemini/antigravity/brain/c77360cc-e824-4c0d-b170-0c52ac0981c5/predictive_dashboard_1772731450101.png",
            "title": "Machine Learning: Mantenimiento Preventivo",
            "desc": "Detecte problemas antes de que la máquina falle en obra.\n- Algoritmos correlacionan datos de consumo (Vales) y rendimiento (Horas Máquina).\n- Alertas sobre anomalías en temperatura y comportamiento del motor.\n- Maximice la disponibilidad de equipos."
        },
        {
            "img": "/Users/klm95441/.gemini/antigravity/brain/c77360cc-e824-4c0d-b170-0c52ac0981c5/ai_contracts_1772731464591.png",
            "title": "LLMs para Análisis de Contratos",
            "desc": "Mitigación de riesgos legales y cumplimiento de tarifas.\n- Extracción de cláusulas (tarifas horarias, mensuales, días mínimos).\n- Comparación automática contra las normativas de ARAMSA y formato CORP-GEM-F-001.\n- Acelera la Firma de Representantes."
        }
    ]
    
    for item in images:
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add Image
        try:
            pic_path = item["img"]
            if os.path.exists(pic_path):
                pic = slide.shapes.add_picture(pic_path, Inches(0.5), Inches(1.5), height=Inches(5))
        except Exception as e:
            print(f"Error adding image {item['img']}: {e}")
            pass # Ignore if image missing
            
        # Add Title Box
        txBox_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        tf = txBox_title.text_frame
        p = tf.paragraphs[0]
        p.text = item["title"]
        p.font.bold = True
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
        
        # Add Description Box
        txBox_desc = slide.shapes.add_textbox(Inches(5.6), Inches(1.5), Inches(4.3), Inches(4.5))
        txBox_desc.text_frame.word_wrap = True
        txBox_desc.text_frame.text = item["desc"]
        for paragraph in txBox_desc.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.space_before = Pt(14)
            paragraph.font.color.rgb = RGBColor(50, 50, 50)
            
    # Slide 6: Optimización de Operaciones (OR)
    slide = prs.slides.add_slide(blank_slide_layout)
    txBox_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = txBox_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Investigación de Operaciones - Asignación Óptima"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    txBox_body = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    txBox_body.text_frame.word_wrap = True
    body_text = (
        "Algoritmos de Optimización de Flotas:\n\n"
        "- Asignación dinámica de maquinaria propia vs. terceros para distintos proyectos.\n"
        "- Minimización de costos logísticos y de desmovilización.\n"
        "- Reducción del margen de equipo subutilizado (horas mínimas no cumplidas).\n"
        "- Planificación automática alineada al cronograma de la obra.\n\n"
        "Mejorando de forma radical la Valorización de Equipo (CORP-GEM-P-002)."
    )
    txBox_body.text_frame.text = body_text
    for paragraph in txBox_body.text_frame.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.space_before = Pt(10)

    # Save
    prs.save("Propuestas_AI_Bitcorp_ERP.pptx")
    print("Presentation saved to Propuestas_AI_Bitcorp_ERP.pptx")

if __name__ == "__main__":
    create_presentation()
